import os
import math
import re  # 新增：引入正则表达式库
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

class PPTBuilder:
    def __init__(self, output_path, width_inch=13.33, height_inch=7.5):
        """
        初始化 PPT 构建器
        :param output_path: 输出文件路径 (.pptx)
        :param width_inch: PPT 宽度 (英寸), 默认 16:9
        :param height_inch: PPT 高度 (英寸)
        """
        self.output_path = output_path
        self.prs = Presentation()
        self.prs.slide_width = Inches(width_inch)
        self.prs.slide_height = Inches(height_inch)
        
        # 布局常量 (可根据需要调整)
        self.MARGIN_X = Inches(0.1)      # 左右边距
        self.MARGIN_BOTTOM = Inches(0.1) # 底部边距
        self.TITLE_TOP = Inches(0.05)    # 标题顶部位置
        self.TITLE_HEIGHT = Inches(0.5)  # 标题高度
        self.CONTENT_TOP = Inches(0.65)  # 图片内容起始高度 (避开标题)
        self.IMG_GAP = Inches(0.05)      # 图片间隙

    def _add_slide_title(self, slide, text, font_size=20):
        """添加幻灯片标题"""
        title_box = slide.shapes.add_textbox(
            Inches(0.5), self.TITLE_TOP, 
            self.prs.slide_width - Inches(1.0), self.TITLE_HEIGHT
        )
        tf = title_box.text_frame
        tf.text = text
        p = tf.paragraphs[0]
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT # 或 PP_ALIGN.CENTER

    def add_section_cover(self, section_title):
        """添加章节过渡页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6]) # 空白页
        textbox = slide.shapes.add_textbox(
            Inches(0), self.prs.slide_height/2 - Inches(1), 
            self.prs.slide_width, Inches(2)
        )
        tf = textbox.text_frame
        tf.text = section_title
        p = tf.paragraphs[0]
        p.font.size = Pt(44)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    def _place_image_in_slot(self, slide, img_path, left, top, slot_w, slot_h):
        """
        核心绘图函数：在给定插槽内居中绘制图片，保持宽高比
        """
        try:
            with Image.open(img_path) as img:
                img_w, img_h = img.size
                aspect = img_w / img_h

            # 尝试以“宽度”为基准适应
            final_h = slot_w / aspect
            final_w = slot_w

            # 如果高度超出了插槽，则改用“高度”为基准
            if final_h > slot_h:
                final_h = slot_h
                final_w = final_h * aspect

            # 计算居中偏移量
            offset_x = (slot_w - final_w) / 2
            offset_y = (slot_h - final_h) / 2

            slide.shapes.add_picture(
                img_path, 
                left + offset_x, 
                top + offset_y, 
                width=final_w, 
                height=final_h
            )
        except Exception as e:
            print(f"[Error] 图片加载失败: {os.path.basename(img_path)} -> {e}")

    def add_grid_slide(self, title, image_paths, cols=4, rows=3):
        """
        添加一张网格布局的幻灯片
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_slide_title(slide, title)

        if not image_paths:
            return

        # 计算可用绘图区域
        avail_w = self.prs.slide_width - 2 * self.MARGIN_X
        avail_h = self.prs.slide_height - self.CONTENT_TOP - self.MARGIN_BOTTOM
        
        # 计算单个插槽尺寸
        slot_w = (avail_w - (cols - 1) * self.IMG_GAP) / cols
        slot_h = (avail_h - (rows - 1) * self.IMG_GAP) / rows

        for idx, img_path in enumerate(image_paths):
            # 防止图片数量超过网格
            if idx >= cols * rows:
                break
                
            r = idx // cols
            c = idx % cols
            
            x = self.MARGIN_X + c * (slot_w + self.IMG_GAP)
            y = self.CONTENT_TOP + r * (slot_h + self.IMG_GAP)
            
            self._place_image_in_slot(slide, img_path, x, y, slot_w, slot_h)

    def search_images(self, base_dir, sub_folder, patterns=None, model=None, suffix=".png", use_regex=True, model_delimiter="_"):
        """
        通用搜图逻辑 (完全体)
        :param base_dir: 根目录
        :param sub_folder: 目标子文件夹
        :param patterns: 匹配模式列表 (OR 关系)
        :param model: 模型名称
        :param suffix: 文件后缀
        :param use_regex: 是否将 patterns 视为正则表达式
        :param model_delimiter: 模型名与后续字符的连接符（如 "_", "-", 甚至是 ""）
        """
        found_images = []
        if patterns is None:
            patterns = []
            
        search_dirs = []
        full_sub_path = os.path.join(base_dir, sub_folder)
        
        # 策略：如果有 model，尝试进入模型子目录；否则直接搜 sub_folder
        if model:
            p = os.path.join(full_sub_path, model)
            if os.path.exists(p):
                search_dirs.append(p)
            elif os.path.exists(full_sub_path):
                 search_dirs.append(full_sub_path)
        else:
            if os.path.exists(full_sub_path):
                search_dirs.append(full_sub_path)

        # 去重
        search_dirs = list(set(search_dirs))

        for d in search_dirs:
            if not os.path.exists(d): continue
            
            files = sorted(os.listdir(d))
            for f in files:
                if not f.endswith(suffix): continue
                
                # 【核心逻辑1】：解决重复索引。利用外部传入的 model_delimiter 组合成前缀匹配
                if model and not f.startswith(f"{model}{model_delimiter}"):
                    continue
                
                # 如果没有设定 patterns，并且走到这一步说明模型前缀匹配成功，直接添加
                if not patterns:
                    found_images.append(os.path.join(d, f))
                    continue
                
                # 【核心逻辑2】：利用 OR 逻辑和正则表达式匹配关键字
                is_match = False
                for p_str in patterns:
                    if use_regex:
                        if re.search(p_str, f):
                            is_match = True
                            break
                    else:
                        if p_str in f:
                            is_match = True
                            break
                
                if is_match:
                    found_images.append(os.path.join(d, f))
                    
        return found_images

    def run_tasks(self, tasks):
        """
        执行任务列表生成 PPT
        """
        print(f"🚀 开始生成 PPT: {self.output_path}")
        
        for task in tasks:
            section_name = task.get('section_name')
            slide_title_prefix = task.get('slide_title', section_name)
            
            # 1. 章节过渡页
            if section_name:
                self.add_section_cover(section_name)
                print(f"\n--- 处理章节: {section_name} ---")

            # 2. 获取配置（暴露核心接口给外部 tasks 配置）
            base_dir = task['base_dir']
            folder = task['folder']
            patterns = task.get('patterns', []) 
            models = task.get('models', [])     
            use_regex = task.get('use_regex', True)             # 默认开启正则
            model_delimiter = task.get('model_delimiter', "_")  # 默认使用 "_" 作为分隔符
            cols = task.get('cols', 2)
            rows = task.get('rows', 1)
            
            # 3. 收集图片
            all_images = []
            
            if models:
                for model in models:
                    # 将正则开关和分隔符一起传给搜图核心函数
                    imgs = self.search_images(
                        base_dir, folder, 
                        patterns=patterns, 
                        model=model, 
                        use_regex=use_regex, 
                        model_delimiter=model_delimiter
                    )
                    all_images.extend(imgs)
            else:
                all_images = self.search_images(
                    base_dir, folder, 
                    patterns=patterns, 
                    use_regex=use_regex, 
                    model_delimiter=model_delimiter
                )

            if not all_images:
                print(f"   [警告] 未找到图片，跳过。条件 -> Folder: {folder}, Patterns: {patterns}")
                continue

            # 4. 分页绘制
            imgs_per_slide = cols * rows
            total_slides = math.ceil(len(all_images) / imgs_per_slide)
            
            for i in range(total_slides):
                start = i * imgs_per_slide
                end = start + imgs_per_slide
                batch = all_images[start:end]
                
                # 标题带页码
                page_title = f"{slide_title_prefix} ({i+1}/{total_slides})" if total_slides > 1 else slide_title_prefix
                
                self.add_grid_slide(page_title, batch, cols=cols, rows=rows)
                print(f"   -> 生成幻灯片: {page_title} (包含 {len(batch)} 张图)")

        # 保存
        os.makedirs(os.path.dirname(self.output_path), exist_ok=True)
        self.prs.save(self.output_path)
        print(f"\n✅ PPT 保存成功: {self.output_path}")